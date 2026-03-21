"""Smoke tests for DAZI PowerPoint processing scripts.

Run with:
    cd /Users/wyatt/Downloads/DAZI/agent-operator
    python3 -m pytest apps/electron/resources/scripts/tests/test_pptx_tool_smoke.py -v

Tests verify the pptx skill scripts under SKILLs/pptx/scripts/.
"""

from __future__ import annotations

import json
import tempfile
import unittest
import zipfile
from pathlib import Path

try:
    from ._tool_test_harness import (
        SKILLS_DIR,
        PPTX_SCRIPTS,
        build_env,
        has_dependency,
        run_script,
    )
except ImportError:
    from _tool_test_harness import (
        SKILLS_DIR,
        PPTX_SCRIPTS,
        build_env,
        has_dependency,
        run_script,
    )


SKIP_REASON = None
if not has_dependency("pptx"):
    SKIP_REASON = "python-pptx not installed"


@unittest.skipIf(SKIP_REASON, SKIP_REASON or "")
class PptxToolSmokeTests(unittest.TestCase):
    """Smoke tests for PowerPoint processing scripts in DAZI."""

    @classmethod
    def setUpClass(cls) -> None:
        cls.env = build_env()
        cls.tmpdir_obj = tempfile.TemporaryDirectory(prefix="dazi-pptx-smoke-")
        cls.tmpdir = Path(cls.tmpdir_obj.name)

    @classmethod
    def tearDownClass(cls) -> None:
        cls.tmpdir_obj.cleanup()

    def _create_test_pptx(self, path: Path, num_slides: int = 3) -> None:
        """Create a test presentation with the given number of slides."""
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        for i in range(num_slides):
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = slide.shapes.title
            if title:
                title.text = f"Slide {i + 1} Title"
            # Add a text box with body content
            from pptx.util import Pt
            txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(2))
            tf = txBox.text_frame
            tf.text = f"Content for slide {i + 1}"
        prs.save(str(path))

    def test_create_and_read_presentation(self) -> None:
        """Verify creating a presentation and reading it back."""
        from pptx import Presentation

        deck = self.tmpdir / "deck.pptx"
        self._create_test_pptx(deck, num_slides=3)
        self.assertTrue(deck.exists())

        prs = Presentation(str(deck))
        self.assertEqual(len(prs.slides), 3)

    def test_inventory_script(self) -> None:
        """inventory.py should extract text content from a presentation."""
        script = PPTX_SCRIPTS / "inventory.py"
        if not script.exists():
            self.skipTest(f"Script not found: {script}")

        deck = self.tmpdir / "inventory_test.pptx"
        self._create_test_pptx(deck, num_slides=2)

        output_json = self.tmpdir / "inventory.json"
        result = run_script(
            script,
            str(deck),
            str(output_json),
            env=self.env,
            cwd=PPTX_SCRIPTS,  # Some scripts import siblings
        )
        self.assertEqual(result.returncode, 0, msg=result.stderr)
        self.assertTrue(output_json.exists(), "Expected inventory.json output")

        data = json.loads(output_json.read_text())
        # inventory.py produces a dict keyed by slide ID
        self.assertIsInstance(data, dict)
        self.assertGreaterEqual(len(data), 1, "Expected at least one slide in inventory")

    def test_rearrange_script(self) -> None:
        """rearrange.py should reorder slides according to a sequence."""
        script = PPTX_SCRIPTS / "rearrange.py"
        if not script.exists():
            self.skipTest(f"Script not found: {script}")

        deck = self.tmpdir / "rearrange_input.pptx"
        self._create_test_pptx(deck, num_slides=3)

        output = self.tmpdir / "rearranged.pptx"
        # Reverse the slides: 2,1,0
        result = run_script(
            script,
            str(deck),
            str(output),
            "2,1,0",
            env=self.env,
            cwd=PPTX_SCRIPTS,
        )
        self.assertEqual(result.returncode, 0, msg=result.stderr)
        self.assertTrue(output.exists())

        from pptx import Presentation
        prs = Presentation(str(output))
        self.assertEqual(len(prs.slides), 3)

    def test_rearrange_invalid_index_fails(self) -> None:
        """rearrange.py should fail gracefully with out-of-range slide index."""
        script = PPTX_SCRIPTS / "rearrange.py"
        if not script.exists():
            self.skipTest(f"Script not found: {script}")

        deck = self.tmpdir / "rearrange_bad.pptx"
        self._create_test_pptx(deck, num_slides=2)

        output = self.tmpdir / "rearranged_bad.pptx"
        result = run_script(
            script,
            str(deck),
            str(output),
            "99",
            env=self.env,
            cwd=PPTX_SCRIPTS,
        )
        self.assertNotEqual(result.returncode, 0, "Expected failure for out-of-range index")

    def test_ooxml_unpack_pptx(self) -> None:
        """ooxml unpack.py should extract XML from a .pptx file."""
        unpack_script = SKILLS_DIR / "pptx" / "ooxml" / "scripts" / "unpack.py"
        if not unpack_script.exists():
            self.skipTest(f"Script not found: {unpack_script}")

        deck = self.tmpdir / "unpack_test.pptx"
        self._create_test_pptx(deck, num_slides=1)

        output_dir = self.tmpdir / "pptx_unpacked"
        result = run_script(
            unpack_script,
            str(deck),
            str(output_dir),
            env=self.env,
        )
        self.assertEqual(result.returncode, 0, msg=result.stderr)
        self.assertTrue(output_dir.exists())

        # Should have ppt/slides/ directory
        slides_dir = output_dir / "ppt" / "slides"
        self.assertTrue(
            slides_dir.exists(),
            f"Expected ppt/slides/ in unpacked output, got: {list(output_dir.rglob('*'))}",
        )

    def test_pptx_is_valid_zip(self) -> None:
        """Verify that .pptx files created by python-pptx are valid ZIPs."""
        deck = self.tmpdir / "zip_check.pptx"
        self._create_test_pptx(deck, num_slides=1)
        self.assertTrue(zipfile.is_zipfile(str(deck)))


if __name__ == "__main__":
    unittest.main()
